#!/usr/bin/env python3
"""
Build VibeLoom pitch deck as a PPTX file.

Mirrors v03/pitch-deck.html (10 slides, 16:9, YC seed-round structure).
PPTX is a portable VC-friendly artifact: opens in PowerPoint, Google Slides,
Keynote. EWOR's deck designer can re-style; VCs can preview in their
inboxes without browsers.

Run:    python3 build-pitch-deck.py
Out:    v03/pitch-deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# --- brand tokens (mirror codæ-manifesto.html) -------------------------------
INK         = RGBColor(0x0A, 0x0A, 0x0A)
INK_2       = RGBColor(0x1A, 0x1A, 0x1A)
INK_3       = RGBColor(0x3A, 0x3A, 0x3A)
INK_4       = RGBColor(0x5A, 0x5A, 0x5A)
INK_5       = RGBColor(0x70, 0x70, 0x70)
LINE        = RGBColor(0xE6, 0xE4, 0xE0)
LINE_2      = RGBColor(0xD6, 0xD3, 0xCD)
BG          = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT     = RGBColor(0xF7, 0xF7, 0xF6)
BG_MUTE     = RGBColor(0xEF, 0xEE, 0xEC)
RED         = RGBColor(0xE8, 0x40, 0x57)
RED_DEEP    = RGBColor(0xC4, 0x3A, 0x50)
RED_TINT    = RGBColor(0xFB, 0xE6, 0xEA)   # solid approximation of rgba

F_SANS  = "Inter"
F_MONO  = "JetBrains Mono"
F_SERIF = "Fraunces"

# --- slide canvas: 13.333" x 7.5" (16:9 widescreen) --------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Generous margins
MARGIN_X = Inches(0.7)
MARGIN_Y_TOP = Inches(0.45)
MARGIN_Y_BOTTOM = Inches(0.45)


# === helpers =================================================================

def add_text(slide, x, y, w, h, text, *, font=F_SANS, size=14, bold=False,
             italic=False, color=INK_2, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.4):
    """Add text with mixed runs. `runs` is a list of dicts:
       {text, font, size, bold, italic, color}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r in runs:
        if r.get("newline"):
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            continue
        run = p.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", F_SANS)
        run.font.size = Pt(r.get("size", 14))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", INK_2)
    return tb


def add_box(slide, x, y, w, h, *, fill=None, line=None, line_width=Pt(0.75),
            line_left_color=None, line_left_width=Pt(2.5), corner=Inches(0.08)):
    """Rounded rectangle background block. Optionally with a thicker red left bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # Set corner radius
    shape.adjustments[0] = 0.06  # ~rounded
    shape.fill.solid()
    if fill is not None:
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    shape.shadow.inherit = False
    # Optional left accent bar (separate thin rectangle on left edge)
    if line_left_color is not None:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, line_left_width, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = line_left_color
        bar.line.fill.background()
        bar.shadow.inherit = False
    return shape


def add_simple_rect(slide, x, y, w, h, *, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_slide_chrome(slide, num, total, eyebrow):
    """Page number + eyebrow at top, footer at bottom."""
    # Top header: slide num (red, mono) + eyebrow (mono)
    add_text(slide, MARGIN_X, MARGIN_Y_TOP,
             Inches(2.5), Inches(0.3),
             f"{num:02d} / {total:02d}",
             font=F_MONO, size=10, bold=True, color=RED_DEEP, line_spacing=1)

    add_text(slide, SLIDE_W - MARGIN_X - Inches(4), MARGIN_Y_TOP,
             Inches(4), Inches(0.3),
             eyebrow.upper(),
             font=F_MONO, size=10, bold=True, color=INK_5,
             align=PP_ALIGN.RIGHT, line_spacing=1)

    # Bottom footer: vibeloom.ai · slide label
    line_y = SLIDE_H - Inches(0.65)
    line = slide.shapes.add_connector(1, MARGIN_X, line_y,
                                       SLIDE_W - MARGIN_X, line_y)
    line.line.color.rgb = LINE
    line.line.width = Pt(0.5)

    add_text(slide, MARGIN_X, SLIDE_H - Inches(0.55),
             Inches(3), Inches(0.3),
             "VIBELOOM.AI",
             font=F_MONO, size=9, bold=True, color=INK_5, line_spacing=1)
    add_text(slide, SLIDE_W - MARGIN_X - Inches(4), SLIDE_H - Inches(0.55),
             Inches(4), Inches(0.3),
             f"{num:02d} — {eyebrow.upper()}",
             font=F_MONO, size=9, bold=True, color=INK_5,
             align=PP_ALIGN.RIGHT, line_spacing=1)


def add_headline(slide, x, y, w, h, parts):
    """Headline with mixed runs (bold sans + italic-serif red accents).
       parts: list of (text, accent_bool)
    """
    runs = []
    for text, accent in parts:
        if accent:
            runs.append({
                "text": text,
                "font": F_SERIF, "size": 30, "italic": True, "bold": False,
                "color": RED,
            })
        else:
            runs.append({
                "text": text,
                "font": F_SANS, "size": 30, "bold": True, "italic": False,
                "color": INK,
            })
    add_runs(slide, x, y, w, h, runs, line_spacing=1.1)


def add_subhead(slide, x, y, w, h, text):
    add_text(slide, x, y, w, h, text,
             font=F_SANS, size=14, color=INK_3, line_spacing=1.4)


# === build deck ==============================================================

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank

    # =========================================================================
    # SLIDE 1 — TITLE
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 1, 10, "Title")

    # Pre-seed badge top-right (red border pill — approximated)
    badge_x = SLIDE_W - MARGIN_X - Inches(3.0)
    add_box(s, badge_x, MARGIN_Y_TOP - Inches(0.05),
            Inches(3.0), Inches(0.36),
            fill=RED_TINT, line=RED, line_width=Pt(0.75))
    add_text(s, badge_x, MARGIN_Y_TOP + Inches(0.02),
             Inches(3.0), Inches(0.3),
             "PRE-SEED THESIS · MAY 2026",
             font=F_MONO, size=9, bold=True, color=RED_DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1)

    # Title mark "VibeLoom" — Vibe in ink, Loom in red
    add_runs(s, MARGIN_X, Inches(2.0), Inches(10), Inches(2.0), [
        {"text": "Vibe", "font": F_SANS, "size": 96, "bold": True, "color": INK},
        {"text": "Loom", "font": F_SANS, "size": 96, "bold": True, "color": RED},
    ], line_spacing=1)

    # Tagline (italic Fraunces with red accent)
    add_runs(s, MARGIN_X, Inches(3.7), Inches(11), Inches(0.7), [
        {"text": "The contract layer for ", "font": F_SERIF, "size": 28, "italic": True, "color": INK_2},
        {"text": "agentic engineering.", "font": F_SERIF, "size": 28, "italic": True, "color": RED},
    ], line_spacing=1.2)

    # One-liner — leads with dark-factory framing
    add_runs(s, MARGIN_X, Inches(4.55), Inches(10.5), Inches(1.4), [
        {"text": "Code becomes a dark factory. Humans operate one level up.",
         "size": 14, "bold": True, "color": INK},
        {"text": " We build the contract substrate that keeps AI-generated code coherent across cycles — engineering teams approve intent and architecture; agents regenerate code from approved contract; drift is detected before it ships.",
         "size": 14, "color": INK_3},
    ], line_spacing=1.5)

    # Founder meta — bottom row (mono)
    meta_y = Inches(6.45)
    add_runs(s, MARGIN_X, meta_y, Inches(12), Inches(0.4), [
        {"text": "Ilya Baimetov", "font": F_MONO, "size": 10, "bold": True, "color": INK},
        {"text": ", Founder      ", "font": F_MONO, "size": 10, "color": INK_4},
        {"text": "ilya.baimetov@vibeloom.ai      ", "font": F_MONO, "size": 10, "color": INK_4},
        {"text": "vibeloom.ai      ", "font": F_MONO, "size": 10, "color": INK_4},
        {"text": "May 2026 · v0.3", "font": F_MONO, "size": 10, "color": INK_4},
    ], line_spacing=1)

    # =========================================================================
    # SLIDE 2 — PROBLEM
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 2, 10, "Problem")

    add_headline(s, MARGIN_X, Inches(0.95), Inches(12), Inches(1.4), [
        ("AI is shipping code ", False),
        ("faster", True),
        (" than anyone can govern it.", False),
    ])

    add_runs(s, MARGIN_X, Inches(2.15), Inches(11.5), Inches(0.85), [
        {"text": "Engineering teams running multi-agent code-gen in production are paying a measurable tax. Existing tooling — Cursor, Copilot, Spec Kit — optimizes the moment of generation. ",
         "size": 13, "color": INK_3},
        {"text": "The drift between cycles is unsolved.", "size": 13, "bold": True, "color": INK},
    ], line_spacing=1.5)

    # Hero number 22.7%
    hero_y = Inches(3.2)
    add_runs(s, MARGIN_X, hero_y, Inches(3.5), Inches(1.4), [
        {"text": "22.7%", "font": F_SERIF, "size": 64, "italic": True, "bold": True, "color": RED},
    ], line_spacing=0.95)

    add_runs(s, Inches(4.5), hero_y + Inches(0.18), Inches(8), Inches(1.3), [
        {"text": "of ", "size": 14, "color": INK_2},
        {"text": "tracked AI-introduced issues survive", "size": 14, "bold": True, "color": INK},
        {"text": " at the latest revision — across ", "size": 14, "color": INK_2},
        {"text": "302.6K commits", "size": 14, "bold": True, "color": INK},
        {"text": " in ", "size": 14, "color": INK_2},
        {"text": "6,299 production repos", "size": 14, "bold": True, "color": INK},
        {"text": ".\n", "size": 14, "color": INK_2},
        {"text": "LIU ET AL., \"AI-DEBT IN THE WILD\" · MAR 2026", "font": F_MONO, "size": 9, "color": INK_5},
    ], line_spacing=1.4)

    # 3 supporting stats
    cards_y = Inches(4.95)
    card_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) / 3
    cards = [
        ("3–5×",
         "short-term velocity gain from Cursor — ",
         "dissipates in 2 months,",
         " leaving +30% warnings, +41% complexity",
         "CMU · MSR 2026"),
        ("17%",
         "lower comprehension scores for AI-assisted developers — ",
         "largest decline in debugging",
         "",
         "Anthropic · O'Reilly · Apr 2026"),
        ("60% / 0–20%",
         "share of work that uses AI / share of tasks teams can ",
         "fully delegate",
         " without supervision",
         "Anthropic · Q1 2026"),
    ]
    for i, (num, lbl_a, lbl_b, lbl_c, src) in enumerate(cards):
        cx = MARGIN_X + i * (card_w + Inches(0.2))
        add_box(s, cx, cards_y, card_w, Inches(1.7),
                fill=BG_SOFT, line_left_color=RED)
        add_text(s, cx + Inches(0.25), cards_y + Inches(0.1),
                 card_w - Inches(0.4), Inches(0.5),
                 num, font=F_SERIF, size=22, italic=True, bold=True, color=INK,
                 line_spacing=1.1)
        add_runs(s, cx + Inches(0.25), cards_y + Inches(0.55),
                 card_w - Inches(0.4), Inches(0.85), [
                    {"text": lbl_a, "size": 10, "color": INK_3},
                    {"text": lbl_b, "size": 10, "bold": True, "color": INK},
                    {"text": lbl_c, "size": 10, "color": INK_3},
                 ], line_spacing=1.4)
        add_text(s, cx + Inches(0.25), cards_y + Inches(1.45),
                 card_w - Inches(0.4), Inches(0.25),
                 src.upper(), font=F_MONO, size=8, color=RED_DEEP, line_spacing=1)

    # =========================================================================
    # SLIDE 3 — THE BET (dark factory)
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 3, 10, "The bet")

    add_headline(s, MARGIN_X, Inches(0.95), Inches(12), Inches(1.4), [
        ("Code becomes a ", False),
        ("dark factory.", True),
        (" We build the ", False),
        ("contract layer", True),
        (" above it.", False),
    ])

    add_subhead(s, MARGIN_X, Inches(2.55), Inches(11.5), Inches(0.85),
                "Lights-out coding. Humans approve intent, product, and architecture. Our deterministic engine regenerates code from the approved contract every cycle. Code is machinery, not literature — generated, regenerated, never maintained by hand.")

    # Before / Arrow / After
    ba_y = Inches(3.6)
    ba_h = Inches(2.6)
    col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.6)) / 2 - Inches(0.2)

    # Before
    bx = MARGIN_X
    add_box(s, bx, ba_y, col_w, ba_h, fill=BG_SOFT, line=LINE, line_width=Pt(0.75))
    add_text(s, bx + Inches(0.3), ba_y + Inches(0.18), col_w - Inches(0.6), Inches(0.3),
             "TODAY · THE CURSOR ERA", font=F_MONO, size=10, bold=True, color=INK_4, line_spacing=1)
    add_text(s, bx + Inches(0.3), ba_y + Inches(0.55), col_w - Inches(0.6), Inches(0.5),
             "Humans maintain code.", font=F_SANS, size=18, bold=True, color=INK, line_spacing=1.2)
    add_text(s, bx + Inches(0.3), ba_y + Inches(1.15), col_w - Inches(0.6), Inches(0.6),
             "~100,000 LOC", font=F_SERIF, size=22, italic=True, bold=True, color=INK_2, line_spacing=1.1)
    add_text(s, bx + Inches(0.3), ba_y + Inches(1.85), col_w - Inches(0.6), Inches(0.8),
             "Every cycle, every contributor, every agent. Drift compounds invisibly. 22.7% of issues ship.",
             font=F_SANS, size=11, color=INK_3, line_spacing=1.45)

    # Arrow
    arrow_x = bx + col_w + Inches(0.05)
    add_text(s, arrow_x, ba_y + Inches(1.0), Inches(0.5), Inches(0.6),
             "→", font=F_SERIF, size=30, italic=True, bold=True, color=RED,
             align=PP_ALIGN.CENTER, line_spacing=1)

    # After (dark factory)
    ax = bx + col_w + Inches(0.6)
    add_box(s, ax, ba_y, col_w, ba_h, fill=RED_TINT, line=RED, line_width=Pt(0.75),
            line_left_color=RED)
    add_text(s, ax + Inches(0.3), ba_y + Inches(0.18), col_w - Inches(0.6), Inches(0.3),
             "AGENTIC ENGINEERING · DARK FACTORY",
             font=F_MONO, size=10, bold=True, color=RED_DEEP, line_spacing=1)
    add_text(s, ax + Inches(0.3), ba_y + Inches(0.55), col_w - Inches(0.6), Inches(0.5),
             "Humans approve contract; factory ships code.",
             font=F_SANS, size=16, bold=True, color=INK, line_spacing=1.2)
    add_text(s, ax + Inches(0.3), ba_y + Inches(1.15), col_w - Inches(0.6), Inches(0.6),
             "~30 contract items", font=F_SERIF, size=22, italic=True, bold=True, color=RED, line_spacing=1.1)
    add_text(s, ax + Inches(0.3), ba_y + Inches(1.85), col_w - Inches(0.6), Inches(0.8),
             "One artifact governs many cycles. Code regenerates deterministically. Drift is detected before ship.",
             font=F_SANS, size=11, color=INK_2, line_spacing=1.45)

    # Tagline at bottom — now includes the dark factory phrase
    add_runs(s, MARGIN_X, Inches(6.45), Inches(12), Inches(0.5), [
        {"text": "Cursor enabled ", "font": F_SERIF, "size": 14, "italic": True, "color": INK_2},
        {"text": "vibe coding.", "font": F_SERIF, "size": 14, "italic": True, "color": RED},
        {"text": " VibeLoom enables ", "font": F_SERIF, "size": 14, "italic": True, "color": INK_2},
        {"text": "agentic engineering", "font": F_SERIF, "size": 14, "italic": True, "color": RED},
        {"text": " — and the ", "font": F_SERIF, "size": 14, "italic": True, "color": INK_2},
        {"text": "dark factory", "font": F_SERIF, "size": 14, "italic": True, "color": RED},
        {"text": " it requires.", "font": F_SERIF, "size": 14, "italic": True, "color": INK_2},
    ], align=PP_ALIGN.CENTER, line_spacing=1.4)

    # =========================================================================
    # SLIDE 4 — WHY NOW
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 4, 10, "Why now")

    add_headline(s, MARGIN_X, Inches(0.95), Inches(12), Inches(1.4), [
        ("Four forces converged in the ", False),
        ("same six-month window.", True),
    ])

    add_subhead(s, MARGIN_X, Inches(2.4), Inches(11.5), Inches(0.6),
                "2026 is the year the contract layer is both buildable and necessary. None of these conditions held two years ago.")

    forces = [
        ("01", "Cursor at $9B proves AI-dev-infra is a real market.",
         "Investors and buyers know AI coding tools are infrastructure now, not experiments. The ", "layer above", " Cursor is the next category."),
        ("02", "Drift is now measurable, not anecdotal.",
         "", "SlopCodeBench", " and \"AI-debt in the wild\" shipped Mar 2026 — drift, erosion, and surviving defects quantified across thousands of real repos."),
        ("03", "Frontier models cross the deterministic-regen threshold.",
         "Claude Opus 4.7, GPT-5, Gemini 3 are reliable enough that regen from approved contracts converges to equivalent output. ", "Two years ago this didn't work.", ""),
        ("04", "The oversight bottleneck is the new constraint.",
         "Anthropic 2026 Trends Report: developers use AI in ", "60% of work, but can fully delegate only 0–20%.", " Buyers know they need a governance layer."),
    ]
    grid_y = Inches(3.25)
    grid_h = Inches(1.62)
    for i, (num, h, pre, bold, post) in enumerate(forces):
        col = i % 2
        row = i // 2
        cx = MARGIN_X + col * (Inches(6.0) + Inches(0.2))
        cy = grid_y + row * (grid_h + Inches(0.2))
        add_box(s, cx, cy, Inches(6.0), grid_h, fill=BG, line=LINE, line_width=Pt(0.6))
        add_text(s, cx + Inches(0.25), cy + Inches(0.18),
                 Inches(0.7), Inches(0.6),
                 num, font=F_SERIF, size=22, italic=True, bold=True, color=RED, line_spacing=1)
        add_text(s, cx + Inches(0.95), cy + Inches(0.18),
                 Inches(4.85), Inches(0.5),
                 h, font=F_SANS, size=12, bold=True, color=INK, line_spacing=1.25)
        add_runs(s, cx + Inches(0.95), cy + Inches(0.78),
                 Inches(4.85), Inches(0.85), [
                    {"text": pre, "size": 10, "color": INK_3},
                    {"text": bold, "size": 10, "bold": True, "color": INK_2},
                    {"text": post, "size": 10, "color": INK_3},
                 ], line_spacing=1.4)

    # =========================================================================
    # SLIDE 5 — PRODUCT
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 5, 10, "Product")

    add_headline(s, MARGIN_X, Inches(0.95), Inches(12), Inches(0.85), [
        ("A methodology, a Skill, and a deterministic engine.", False),
    ])

    add_subhead(s, MARGIN_X, Inches(1.95), Inches(11.5), Inches(0.7),
                "Ships as a Claude Code / Codex Skill on day one — no install. Methodology is open source (MIT). Engine is pure Python with zero runtime dependencies. Five operating modes from solo hacker to enterprise.")

    # Left — flow steps
    flow_x = MARGIN_X
    flow_y = Inches(3.0)
    flow_w = Inches(6.5)
    flow_h = Inches(3.7)

    add_box(s, flow_x, flow_y, flow_w, flow_h, fill=BG_SOFT, line=LINE, line_width=Pt(0.6))

    steps = [
        ("1 · User edits intent.md + approves",
         "capabilities, constraints in plain English; approval trace recorded",
         RED, INK_2),
        ("2 · Engine regenerates downstream",
         "product specs · UX · architecture · code · BDD scenarios",
         INK_3, INK_2),
        ("3 · Eval ladder runs each cycle",
         "decidable structural · mechanical runners · heuristic semantic",
         INK_3, INK_2),
        ("4 · Code ships; code-sync trace closes the loop",
         "every code path traces back to its contract item",
         LINE_2, INK_4),
    ]
    step_h = Inches(0.7)
    step_gap = Inches(0.18)
    for i, (head, sub, accent, txt_color) in enumerate(steps):
        sy = flow_y + Inches(0.25) + i * (step_h + step_gap)
        sx = flow_x + Inches(0.25)
        sw = flow_w - Inches(0.5)
        add_box(s, sx, sy, sw, step_h,
                fill=BG if accent != LINE_2 else BG_MUTE,
                line=LINE, line_width=Pt(0.5),
                line_left_color=accent, line_left_width=Pt(2.5))
        add_text(s, sx + Inches(0.25), sy + Inches(0.06),
                 sw - Inches(0.4), Inches(0.3),
                 head, font=F_MONO, size=11, bold=True, color=txt_color, line_spacing=1.1)
        add_text(s, sx + Inches(0.25), sy + Inches(0.36),
                 sw - Inches(0.4), Inches(0.3),
                 sub, font=F_SANS, size=9, color=INK_4, line_spacing=1.3)

    # Right — modes
    mx = flow_x + flow_w + Inches(0.3)
    mw = SLIDE_W - mx - MARGIN_X
    add_box(s, mx, flow_y, mw, flow_h, fill=BG, line=LINE, line_width=Pt(0.6))

    add_text(s, mx + Inches(0.3), flow_y + Inches(0.2), mw - Inches(0.6), Inches(0.3),
             "FIVE MODES — SAME ENGINE, DIFFERENT SURFACE",
             font=F_MONO, size=9, bold=True, color=RED_DEEP, line_spacing=1)

    modes = [
        ("vibe",   "Solo. Compact stack. No graph. One-way upgrade to full mode."),
        ("pm",     "Product-led. PM owns intent + product specs."),
        ("dev",    "Tech-led. Dev owns intent + system specs."),
        ("ux",     "Design-led. Designer owns intent + UX. Mockups drive product specs."),
        ("expert", "Regulated. Every approval gate explicit. Compliance-ready."),
    ]
    mode_y = flow_y + Inches(0.65)
    for i, (name, desc) in enumerate(modes):
        my = mode_y + i * Inches(0.5)
        add_text(s, mx + Inches(0.3), my, Inches(0.8), Inches(0.4),
                 name, font=F_MONO, size=11, bold=True, color=RED_DEEP, line_spacing=1.2)
        # split desc into bold first word + rest
        first_dot = desc.find(".")
        bold_part = desc[:first_dot+1]
        rest_part = desc[first_dot+1:]
        add_runs(s, mx + Inches(1.15), my, mw - Inches(1.45), Inches(0.45), [
                    {"text": bold_part, "size": 10, "bold": True, "color": INK},
                    {"text": rest_part, "size": 10, "color": INK_3},
                 ], line_spacing=1.35)

    add_text(s, mx + Inches(0.3), flow_y + flow_h - Inches(0.55),
             mw - Inches(0.6), Inches(0.4),
             "Open source: methodology, Skill, templates. Paid: hosted engine, audit, compliance.",
             font=F_MONO, size=9, color=INK_4, line_spacing=1.4)

    # =========================================================================
    # SLIDE 6 — INSIGHT
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 6, 10, "Insight")

    add_headline(s, MARGIN_X, Inches(0.95), Inches(12.2), Inches(1.4), [
        ("Cursor solves ", False),
        ("generation.", True),
        (" We solve ", False),
        ("coherence between generations.", True),
    ])

    add_subhead(s, MARGIN_X, Inches(2.45), Inches(11.5), Inches(0.85),
                "Every existing tool optimizes the moment of code production. The drift problem is structural: it happens between cycles. Solving it requires a different abstraction layer — and a different buyer.")

    # Left: insight statement
    in_x = MARGIN_X
    in_y = Inches(3.6)
    in_w = Inches(5.8)
    in_h = Inches(3.0)
    add_box(s, in_x, in_y, in_w, in_h, fill=BG_SOFT,
            line_left_color=RED, line_left_width=Pt(2.5))
    add_text(s, in_x + Inches(0.3), in_y + Inches(0.2),
             in_w - Inches(0.5), Inches(0.3),
             "THE NON-OBVIOUS BET",
             font=F_MONO, size=9, bold=True, color=RED_DEEP, line_spacing=1)
    add_runs(s, in_x + Inches(0.3), in_y + Inches(0.65),
             in_w - Inches(0.5), in_h - Inches(0.85), [
                {"text": "We are ", "font": F_SERIF, "size": 14, "italic": True, "color": INK},
                {"text": "not in the same fight as Cursor", "font": F_SANS, "size": 14, "bold": True, "color": RED_DEEP},
                {"text": ". Cursor sells productivity to individual developers. We sell ",
                 "font": F_SERIF, "size": 14, "italic": True, "color": INK},
                {"text": "governance", "font": F_SANS, "size": 14, "bold": True, "color": RED_DEEP},
                {"text": " to engineering leaders — at a different layer of the stack, to a different buyer, with a different metric (drift caught vs. lines shipped).",
                 "font": F_SERIF, "size": 14, "italic": True, "color": INK},
             ], line_spacing=1.45)

    # Right: 3 competitor rows
    cmp_x = in_x + in_w + Inches(0.3)
    cmp_w = SLIDE_W - cmp_x - MARGIN_X
    add_text(s, cmp_x, in_y, cmp_w, Inches(0.3),
             "WHERE THE VALUE CHAIN SITS",
             font=F_MONO, size=9, bold=True, color=RED_DEEP, line_spacing=1)

    competitors = [
        ("Cursor / Copilot / Codeium",
         "Vibe coding — chat-driven, in the IDE. Optimizes per-edit velocity.",
         False),
        ("Kiro / Spec Kit / BMAD",
         "Spec-driven — per-feature specs feed agent generation. Specs decay between features.",
         False),
        ("VibeLoom",
         "Contract-driven — system-level contracts govern many cycles. Drift detected at the contract layer, not in code.",
         True),
    ]
    row_y = in_y + Inches(0.45)
    row_h = Inches(0.78)
    row_gap = Inches(0.12)
    for i, (who, what, us) in enumerate(competitors):
        ry = row_y + i * (row_h + row_gap)
        if us:
            add_box(s, cmp_x, ry, cmp_w, row_h, fill=RED_TINT, line=RED, line_width=Pt(0.5))
            who_color = RED_DEEP
        else:
            add_box(s, cmp_x, ry, cmp_w, row_h, fill=BG, line=LINE, line_width=Pt(0.5))
            who_color = INK
        add_text(s, cmp_x + Inches(0.25), ry + Inches(0.1),
                 cmp_w - Inches(0.4), Inches(0.3),
                 who, font=F_MONO, size=11, bold=True, color=who_color, line_spacing=1)
        add_text(s, cmp_x + Inches(0.25), ry + Inches(0.4),
                 cmp_w - Inches(0.4), Inches(0.4),
                 what, font=F_SANS, size=10, color=INK_3 if not us else INK_2,
                 line_spacing=1.4)

    # =========================================================================
    # SLIDE 7 — BUSINESS MODEL
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 7, 10, "Business model")

    add_headline(s, MARGIN_X, Inches(0.95), Inches(12), Inches(1.0), [
        ("Open-source pull. ", False),
        ("Tiered SaaS", True),
        (" capture.", False),
    ])

    add_subhead(s, MARGIN_X, Inches(2.1), Inches(11.5), Inches(0.7),
                "Methodology + Skill open under MIT pulls developers in (PLG flywheel; Cursor's pattern). Paid tiers capture engineering teams that need audit, telemetry, compliance, and on-prem.")

    # 3 pricing cards
    bm_y = Inches(3.1)
    bm_h = Inches(3.2)
    bm_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) / 3
    tiers = [
        ("FREE", "$0", "· MIT", "Solo developers, OSS, hobbyists",
         ["vibe + pm modes", "Methodology + Skill",
          "Self-hosted engine", "100 generations / month"], False),
        ("TEAM · RECOMMENDED", "$30", "/ seat / month", "Engineering teams of 5–50",
         ["Full pm / dev / ux modes", "Hosted engine",
          "Audit logs + drift telemetry", "Unlimited generations",
          "Standard SLA"], True),
        ("ENTERPRISE", "$50K+", "/ year", "Mid-market & enterprise (50+ eng)",
         ["expert mode + compliance", "On-prem engine",
          "SOC2 / HIPAA bundles", "Advisor seat (founder access)",
          "SSO + audit-trail export"], False),
    ]
    for i, (tier, price, unit, who, items, featured) in enumerate(tiers):
        cx = MARGIN_X + i * (bm_w + Inches(0.2))
        if featured:
            add_box(s, cx, bm_y, bm_w, bm_h, fill=RED_TINT, line=RED, line_width=Pt(0.6))
            tier_color = RED_DEEP
            price_color = RED_DEEP
        else:
            add_box(s, cx, bm_y, bm_w, bm_h, fill=BG, line=LINE, line_width=Pt(0.6))
            tier_color = INK_4
            price_color = INK
        add_text(s, cx + Inches(0.3), bm_y + Inches(0.2),
                 bm_w - Inches(0.5), Inches(0.3),
                 tier, font=F_MONO, size=9, bold=True, color=tier_color, line_spacing=1)
        add_runs(s, cx + Inches(0.3), bm_y + Inches(0.55),
                 bm_w - Inches(0.5), Inches(0.7), [
                    {"text": price, "font": F_SERIF, "size": 28, "italic": True, "bold": True, "color": price_color},
                    {"text": " " + unit, "font": F_SANS, "size": 10, "color": INK_4},
                 ], line_spacing=1.1)
        add_text(s, cx + Inches(0.3), bm_y + Inches(1.3),
                 bm_w - Inches(0.5), Inches(0.4),
                 who, font=F_SANS, size=11, bold=True, color=INK_2, line_spacing=1.3)
        for j, item in enumerate(items):
            iy = bm_y + Inches(1.75) + j * Inches(0.28)
            add_text(s, cx + Inches(0.3), iy, Inches(0.2), Inches(0.25),
                     "→", font=F_MONO, size=10, bold=True, color=RED, line_spacing=1)
            add_text(s, cx + Inches(0.55), iy, bm_w - Inches(0.85), Inches(0.25),
                     item, font=F_SANS, size=10, color=INK_3, line_spacing=1.3)

    # GTM line at bottom
    gtm_y = bm_y + bm_h + Inches(0.15)
    add_box(s, MARGIN_X, gtm_y, SLIDE_W - 2 * MARGIN_X, Inches(0.55),
            fill=BG_SOFT, line_left_color=LINE_2, line_left_width=Pt(2))
    add_runs(s, MARGIN_X + Inches(0.25), gtm_y + Inches(0.13),
             SLIDE_W - 2 * MARGIN_X - Inches(0.5), Inches(0.4), [
                {"text": "GTM: ", "size": 10, "bold": True, "color": INK},
                {"text": "land via individual devs (free Skill, Discord, GitHub). Expand to teams when drift becomes a P1. Enterprise via design-partner referrals after 6 months.",
                 "font": F_SERIF, "size": 11, "italic": True, "color": INK_3},
             ], line_spacing=1.4)

    # =========================================================================
    # SLIDE 8 — MARKET
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 8, 10, "Market")

    add_headline(s, MARGIN_X, Inches(0.95), Inches(12.2), Inches(1.0), [
        ("$2.2B SAM at saturation. ", False),
        ("Cursor's $9B", True),
        (" is the proof.", False),
    ])

    add_subhead(s, MARGIN_X, Inches(2.05), Inches(11.5), Inches(0.85),
                "If Cursor's run-rate proves the demand for AI dev infrastructure, the contract layer above it is the natural follow-on category — sold to a different buyer (CTO/VPE) at a different ACV (per-team, not per-seat).")

    # Left: bottom-up calc
    calc_x = MARGIN_X
    calc_y = Inches(3.3)
    calc_w = Inches(6.3)
    calc_h = Inches(3.4)
    add_box(s, calc_x, calc_y, calc_w, calc_h, fill=BG_SOFT, line=LINE, line_width=Pt(0.6))
    add_text(s, calc_x + Inches(0.3), calc_y + Inches(0.2),
             calc_w - Inches(0.5), Inches(0.3),
             "BOTTOM-UP SAM, TRANSPARENT CALC",
             font=F_MONO, size=9, bold=True, color=RED_DEEP, line_spacing=1)

    rows = [
        ("Developers worldwide",      "Stack Overflow Dev Survey 2025",  "~30M"),
        ("Using AI coding tools (33%)", "JetBrains 2025; GitHub Octoverse 2025", "~10M"),
        ("Running multi-cycle agentic generation (30% of those)", "the cohort that hits the drift problem", "~3M"),
        ("Avg seat ACV (mid-tier B2B SaaS)", "", "$720 / yr"),
    ]
    rrow_y = calc_y + Inches(0.65)
    for i, (label, sub, num) in enumerate(rows):
        ry = rrow_y + i * Inches(0.5)
        add_text(s, calc_x + Inches(0.3), ry, calc_w - Inches(2), Inches(0.25),
                 label, font=F_SANS, size=11, color=INK_2, line_spacing=1.3)
        if sub:
            add_text(s, calc_x + Inches(0.3), ry + Inches(0.25), calc_w - Inches(2), Inches(0.2),
                     sub, font=F_MONO, size=8, color=INK_5, line_spacing=1)
        add_text(s, calc_x + calc_w - Inches(1.6), ry, Inches(1.4), Inches(0.3),
                 num, font=F_MONO, size=12, bold=True, color=INK,
                 align=PP_ALIGN.RIGHT, line_spacing=1.2)
        # divider
        if i < len(rows) - 1:
            ly = ry + Inches(0.46)
            line = s.shapes.add_connector(1, calc_x + Inches(0.3), ly,
                                           calc_x + calc_w - Inches(0.3), ly)
            line.line.color.rgb = LINE_2
            line.line.width = Pt(0.5)

    # Total row
    total_y = calc_y + calc_h - Inches(0.65)
    line_t = s.shapes.add_connector(1, calc_x + Inches(0.3), total_y - Inches(0.05),
                                     calc_x + calc_w - Inches(0.3), total_y - Inches(0.05))
    line_t.line.color.rgb = RED
    line_t.line.width = Pt(1.5)
    add_text(s, calc_x + Inches(0.3), total_y, calc_w - Inches(2), Inches(0.4),
             "SAM at saturation", font=F_SANS, size=12, bold=True, color=INK, line_spacing=1.2)
    add_text(s, calc_x + calc_w - Inches(1.6), total_y - Inches(0.05), Inches(1.4), Inches(0.5),
             "$2.2B", font=F_SERIF, size=22, italic=True, bold=True, color=RED,
             align=PP_ALIGN.RIGHT, line_spacing=1)

    # Right: narrative
    nar_x = calc_x + calc_w + Inches(0.3)
    nar_w = SLIDE_W - nar_x - MARGIN_X
    add_text(s, nar_x, calc_y + Inches(0.05), nar_w, Inches(1.0),
             "A new category in AI dev infrastructure.",
             font=F_SERIF, size=18, italic=True, bold=True, color=INK, line_spacing=1.25)

    add_runs(s, nar_x, calc_y + Inches(1.3), nar_w, Inches(1.3), [
                {"text": "Cursor's $9B revenue/valuation is the headline; Copilot, Codeium, Tabnine, Replit fill out the field. ",
                 "size": 11, "color": INK_3},
                {"text": "None of them sells to the buyer we sell to", "size": 11, "bold": True, "color": INK},
                {"text": " — the engineering leader responsible for codebase coherence over time.",
                 "size": 11, "color": INK_3},
            ], line_spacing=1.5)

    # Adjacent
    adj_y = calc_y + Inches(2.5)
    add_box(s, nar_x, adj_y, nar_w, Inches(0.85),
            fill=BG_SOFT, line_left_color=LINE_2, line_left_width=Pt(2))
    add_runs(s, nar_x + Inches(0.25), adj_y + Inches(0.15),
             nar_w - Inches(0.45), Inches(0.7), [
                {"text": "Adjacent expansion: $5B+. ", "size": 10, "bold": True, "color": INK},
                {"text": "Methodology consulting, contract pattern marketplaces, audit + compliance bundles, training. Each unlocked once trace-derived learning ships (year 2).",
                 "size": 10, "color": INK_3},
            ], line_spacing=1.45)

    # =========================================================================
    # SLIDE 9 — TEAM
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 9, 10, "Team")

    add_headline(s, MARGIN_X, Inches(0.95), Inches(12), Inches(1.0), [
        ("Solo founder today. ", False),
        ("Team-of-2", True),
        (" by month six.", False),
    ])

    add_subhead(s, MARGIN_X, Inches(2.05), Inches(11.5), Inches(0.85),
                "VCs commonly flag solo founders. The concern is fair. We treat co-founder acquisition as a tracked Q1 milestone — not glossed over. The methodology, spec, and engine shipped solo are themselves the artifact that recruits the right co-founder.")

    # Left: bio
    bio_x = MARGIN_X
    bio_y = Inches(3.25)
    bio_w = Inches(6.0)
    bio_h = Inches(3.5)
    add_box(s, bio_x, bio_y, bio_w, bio_h, fill=BG, line=LINE, line_width=Pt(0.6))
    add_text(s, bio_x + Inches(0.3), bio_y + Inches(0.2),
             bio_w - Inches(0.5), Inches(0.4),
             "Ilya Baimetov", font=F_SANS, size=18, bold=True, color=INK, line_spacing=1.1)
    add_text(s, bio_x + Inches(0.3), bio_y + Inches(0.62),
             bio_w - Inches(0.5), Inches(0.3),
             "FOUNDER · AUTHOR OF CODÆ & VIBELOOM",
             font=F_MONO, size=9, bold=True, color=RED_DEEP, line_spacing=1)
    add_runs(s, bio_x + Inches(0.3), bio_y + Inches(1.05),
             bio_w - Inches(0.5), Inches(0.7), [
                {"text": "Background. ", "size": 10, "bold": True, "color": INK},
                {"text": "[TODO · Ilya] ", "font": F_MONO, "size": 9, "bold": True, "color": RED_DEEP},
                {"text": "Prior roles, technical depth, shipped products. EWOR weights this heavily — fill before submission.",
                 "size": 10, "color": INK_3},
            ], line_spacing=1.4)
    add_runs(s, bio_x + Inches(0.3), bio_y + Inches(1.85),
             bio_w - Inches(0.5), Inches(1.0), [
                {"text": "Why me. ", "size": 10, "bold": True, "color": INK},
                {"text": "Authored codæ paradigm and v0.3 spec end-to-end: methodology, implementation, comparison whitepaper, 35 generation-ready templates, marketing site. Built daily with frontier agents through 2025–2026; saw the slop pattern emerge in real codebases.",
                 "size": 10, "color": INK_3},
            ], line_spacing=1.4)
    add_runs(s, bio_x + Inches(0.3), bio_y + Inches(2.85),
             bio_w - Inches(0.5), Inches(0.6), [
                {"text": "Insight. ", "size": 10, "bold": True, "color": INK},
                {"text": "Cursor proved chat-driven coding works at scale. The next layer is making the contract — not the chat — the durable surface humans operate on. ",
                 "size": 10, "color": INK_3},
                {"text": "Build the layer above.", "size": 10, "bold": True, "color": INK},
            ], line_spacing=1.4)

    # Right: solo callout + advisors
    sx = bio_x + bio_w + Inches(0.3)
    sw = SLIDE_W - sx - MARGIN_X

    # Solo callout
    so_h = Inches(1.7)
    add_box(s, sx, bio_y, sw, so_h, fill=BG_SOFT,
            line_left_color=INK_4, line_left_width=Pt(2.5))
    add_text(s, sx + Inches(0.3), bio_y + Inches(0.2), sw - Inches(0.5), Inches(0.3),
             "ON BEING SOLO (THE ELEPHANT)",
             font=F_MONO, size=9, bold=True, color=INK_4, line_spacing=1)
    add_runs(s, sx + Inches(0.3), bio_y + Inches(0.6), sw - Inches(0.5), so_h - Inches(0.7), [
                {"text": "Q1 milestone, tracked: ", "font": F_SANS, "size": 11, "bold": True, "color": INK},
                {"text": "technical co-founder onboarded by month 6. Profile: distributed-systems / dev-infra / ex-platform-engineering. The pre-seed round itself is the team-formation forcing function — not a workaround.",
                 "font": F_SERIF, "size": 11, "italic": True, "color": INK_2},
            ], line_spacing=1.45)

    # Advisors
    av_y = bio_y + so_h + Inches(0.15)
    av_h = bio_h - so_h - Inches(0.15)
    add_box(s, sx, av_y, sw, av_h, fill=BG, line=LINE, line_width=Pt(0.6))
    add_text(s, sx + Inches(0.3), av_y + Inches(0.2), sw - Inches(0.5), Inches(0.3),
             "ADVISORS & PIPELINE",
             font=F_MONO, size=9, bold=True, color=RED_DEEP, line_spacing=1)
    add_runs(s, sx + Inches(0.3), av_y + Inches(0.6), sw - Inches(0.5), av_h - Inches(0.7), [
                {"text": "[TODO · Ilya] ", "font": F_MONO, "size": 9, "bold": True, "color": RED_DEEP},
                {"text": "Named advisor list (or \"advisory list in formation, X conversations active\"). Plus: design-partner pipeline of N teams currently scoping pilots.",
                 "size": 11, "color": INK_3},
            ], line_spacing=1.45)

    # =========================================================================
    # SLIDE 10 — ASK (VC-independent)
    # =========================================================================
    s = prs.slides.add_slide(blank_layout)
    add_slide_chrome(s, 10, 10, "Ask")

    add_headline(s, MARGIN_X, Inches(0.95), Inches(10.5), Inches(1.0), [
        ("Pre-seed: ", False),
        ("$XXX", True),
        (" for 18 months of runway.", False),
    ])

    # TODO badge inline near the headline
    todo_x = MARGIN_X + Inches(10.6)
    add_box(s, todo_x, Inches(1.05), Inches(2.0), Inches(0.32),
            fill=RED_TINT, line=RED, line_width=Pt(0.6))
    add_text(s, todo_x, Inches(1.1),
             Inches(2.0), Inches(0.3),
             "TODO · ILYA: TARGET $",
             font=F_MONO, size=8, bold=True, color=RED_DEEP,
             align=PP_ALIGN.CENTER, line_spacing=1)

    add_subhead(s, MARGIN_X, Inches(2.1), Inches(11.5), Inches(0.85),
                "Build the contract layer for the agentic era. Engine v0.4 → first 10 design partners → drift telemetry → Series A traction. Solo today; team-of-2 by month six.")

    # 2 cards
    ask_y = Inches(3.35)
    ask_h = Inches(2.4)
    ask_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.3)) / 2

    # Primary (red-tinted) — Use of Funds
    add_box(s, MARGIN_X, ask_y, ask_w, ask_h, fill=RED_TINT, line=RED, line_width=Pt(0.6),
            line_left_color=RED, line_left_width=Pt(2.5))
    add_text(s, MARGIN_X + Inches(0.3), ask_y + Inches(0.2),
             ask_w - Inches(0.5), Inches(0.3),
             "USE OF FUNDS — 18 MONTHS",
             font=F_MONO, size=9, bold=True, color=RED_DEEP, line_spacing=1)
    funds_items = [
        ("Engineering — 55%", " · technical co-founder + 2 senior engineers (engine, infra)"),
        ("Design partners — 25%", " · dedicated success + onboarding + drift-telemetry tooling"),
        ("Legal / IP / compliance — 10%", " · entity, IP assignments, SOC2 prep"),
        ("Buffer + ops — 10%", " · cloud, tooling, contingency"),
    ]
    for i, (bold_part, rest) in enumerate(funds_items):
        iy = ask_y + Inches(0.65) + i * Inches(0.42)
        add_text(s, MARGIN_X + Inches(0.3), iy, Inches(0.2), Inches(0.3),
                 "→", font=F_MONO, size=11, bold=True, color=RED, line_spacing=1)
        add_runs(s, MARGIN_X + Inches(0.55), iy, ask_w - Inches(0.85), Inches(0.4), [
                    {"text": bold_part, "size": 10, "bold": True, "color": INK},
                    {"text": rest, "size": 10, "color": INK_2},
                 ], line_spacing=1.35)

    # Right (white) — Milestones to Series A
    rx = MARGIN_X + ask_w + Inches(0.3)
    add_box(s, rx, ask_y, ask_w, ask_h, fill=BG, line=LINE, line_width=Pt(0.6))
    add_text(s, rx + Inches(0.3), ask_y + Inches(0.2),
             ask_w - Inches(0.5), Inches(0.3),
             "MILESTONES TO SERIES A",
             font=F_MONO, size=9, bold=True, color=RED_DEEP, line_spacing=1)
    milestones = [
        ("Month 3", " · Engine v0.4 dogfood-ready (spec → runnable Python)"),
        ("Month 6", " · Technical co-founder onboarded + 5 design partners shipping"),
        ("Month 12", " · First paying teams + drift telemetry across 10+ codebases"),
        ("Month 18", " · Series A on the back of trace-derived learning evidence"),
    ]
    for i, (bold_part, rest) in enumerate(milestones):
        iy = ask_y + Inches(0.65) + i * Inches(0.42)
        add_text(s, rx + Inches(0.3), iy, Inches(0.2), Inches(0.3),
                 "→", font=F_MONO, size=11, bold=True, color=RED, line_spacing=1)
        add_runs(s, rx + Inches(0.55), iy, ask_w - Inches(0.85), Inches(0.4), [
                    {"text": bold_part, "size": 10, "bold": True, "color": INK},
                    {"text": rest, "size": 10, "color": INK_2},
                 ], line_spacing=1.35)

    # Closer (dark band)
    close_y = ask_y + ask_h + Inches(0.2)
    close_h = Inches(1.0)
    add_simple_rect(s, MARGIN_X, close_y, SLIDE_W - 2 * MARGIN_X, close_h, fill=INK)

    # Tag (left) — now mentions dark factory
    add_runs(s, MARGIN_X + Inches(0.4), close_y + Inches(0.18),
             Inches(7.5), close_h - Inches(0.3), [
                {"text": "Cursor enabled ", "size": 14, "bold": True, "color": BG},
                {"text": "vibe coding.", "font": F_SERIF, "size": 14, "italic": True, "bold": True, "color": RED},
                {"text": "\nVibeLoom enables ", "size": 14, "bold": True, "color": BG},
                {"text": "agentic engineering", "font": F_SERIF, "size": 14, "italic": True, "bold": True, "color": RED},
                {"text": " — and the ", "size": 14, "bold": True, "color": BG},
                {"text": "dark factory", "font": F_SERIF, "size": 14, "italic": True, "bold": True, "color": RED},
                {"text": " it requires.", "size": 14, "bold": True, "color": BG},
            ], line_spacing=1.3)

    # Contact (right)
    ctx_x = MARGIN_X + Inches(8.0)
    ctx_w = SLIDE_W - 2 * MARGIN_X - Inches(8.0) - Inches(0.4)
    add_runs(s, ctx_x, close_y + Inches(0.18),
             ctx_w, close_h - Inches(0.3), [
                {"text": "Ilya Baimetov\n", "font": F_MONO, "size": 10, "bold": True, "color": BG},
                {"text": "ilya.baimetov@vibeloom.ai\n", "font": F_MONO, "size": 9, "color": BG},
                {"text": "vibeloom.ai · github.com/ilya-baimetov/vibeloom", "font": F_MONO, "size": 9, "color": BG},
            ], align=PP_ALIGN.RIGHT, line_spacing=1.5)

    # =========================================================================
    # SAVE
    # =========================================================================
    out_path = "v03/pitch-deck.pptx"
    prs.save(out_path)
    print(f"✓ wrote {out_path}")
    print(f"  slides: {len(prs.slides)}")
    print(f"  size:   {prs.slide_width / 914400:.3f}\" × {prs.slide_height / 914400:.3f}\" (16:9)")


if __name__ == "__main__":
    build()
